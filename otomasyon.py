import streamlit as st
import os
import glob
import copy
import re
from io import BytesIO
from pypdf import PdfReader
from pptx import Presentation

# --- SAYFA AYARLARI ---
st.set_page_config(page_title="MSEP Sunum Hazırlayıcı", page_icon="🎓", layout="centered")

st.title("🎓 MSEP Otomatik Sunum Sistemi")
st.info("Bu sistem GitHub deposundaki güncel PDF ve Şablonu kullanarak sunum hazırlar.")

# --- YARDIMCI FONKSİYONLAR (SENİN MOTORUN) ---
def replace_text_preserve_style(paragraph, new_text):
    if not paragraph.runs:
        paragraph.add_run().text = new_text
        return
    paragraph.runs[0].text = new_text
    for i in range(1, len(paragraph.runs)):
        paragraph.runs[i].text = ""

def dosyalari_bul():
    # GitHub deposundaki (sunucudaki) dosyaları bulur
    pdf_files = glob.glob("*.pdf")
    # pptx bulurken ~$ ile başlayan geçici dosyaları ele
    all_pptx = glob.glob("*.pptx")
    pptx_files = [f for f in all_pptx if not os.path.basename(f).startswith("~$")]
    
    # Şablonu bul (ismi 'sablon' olanı tercih et)
    sablon_file = None
    if pptx_files:
        sablon_file = next((f for f in pptx_files if "sablon" in os.path.basename(f).lower()), pptx_files[0])

    return pdf_files, sablon_file

def tr_lower(text):
    degisim = {"I": "ı", "İ": "i", "Ş": "ş", "Ğ": "ğ", "Ü": "ü", "Ö": "ö", "Ç": "ç"}
    for kaynak, hedef in degisim.items():
        text = text.replace(kaynak, hedef)
    return text.lower()

def verileri_ayikla(pdf_path, secilen_ay):
    text = ""
    try:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    except:
        return []

    lines = [line.strip() for line in text.split('\n') if line.strip()]
    events = []
    current_event = None
    
    date_pattern = re.compile(r"^\d{1,2}\s+[A-ZİıĞÜŞÖÇça-zğüşıöç]+\s+\d{4}\s+[A-ZİıĞÜŞÖÇça-zğüşıöç]+")
    chair_pattern = re.compile(r"^Başkan(?:lar)?\s*:", re.IGNORECASE)
    speaker_pattern = re.compile(r"^(Konuşmacı(?:lar)?)\s*:", re.IGNORECASE)
    
    valid_titles = ["Prof", "Doç", "Uzm", "Dr", "Arş", "Öğr"]
    yasakli_kelimeler = ["Anabilim", "Bilim", "Kurulu", "Programı", "Saat", "Zoom", "Tasarım", "Baskı", "Cihan", "Balcı", "Vaka Sunumları", "Başkanı", "Müdürü", "Rektörü", "Dekanı"]

    i = 0
    while i < len(lines):
        line = lines[i]
        
        if date_pattern.match(line):
            if current_event and current_event['topic'] and (current_event['chair'] or current_event['speakers']):
                events.append(current_event)
            current_event = {'date': line, 'topic': "", 'chair': "", 'speakers': [], 'temp_topic': []}
            i += 1
            continue
        
        if not current_event:
            i += 1
            continue

        if chair_pattern.match(line):
            chair_name = re.sub(r"^Başkan(?:lar)?\s*:\s*", "", line, flags=re.IGNORECASE).strip()
            if not chair_name and i + 1 < len(lines):
                chair_name = lines[i+1].strip()
                i += 1
            current_event['chair'] = chair_name
            if current_event['temp_topic']:
                current_event['topic'] = " ".join(current_event['temp_topic']).strip().strip('“').strip('”').strip('"')
                current_event['temp_topic'] = []
            i += 1
            continue

        speaker_match = speaker_pattern.match(line)
        if speaker_match:
            etiket = speaker_match.group(1).lower()
            is_plural = "lar" in etiket
            first_speaker = re.sub(r"^Konuşmacı(?:lar)?\s*:\s*", "", line, flags=re.IGNORECASE).strip()
            
            def is_valid(s):
                if len(s) < 5 or any(y in s for y in yasakli_kelimeler) or not any(u in s for u in valid_titles) or "2025" in s: return False
                return True

            if first_speaker and is_valid(first_speaker):
                current_event['speakers'].append(first_speaker)

            if not is_plural:
                if current_event['speakers']: 
                    i += 1; continue
                else:
                    if i + 1 < len(lines):
                        nxt = lines[i+1].strip()
                        if is_valid(nxt): current_event['speakers'].append(nxt)
                        i += 2; continue
            else:
                j = i + 1
                while j < len(lines):
                    nxt = lines[j]
                    if date_pattern.match(nxt) or chair_pattern.match(nxt) or speaker_pattern.match(nxt): break
                    if is_valid(nxt.strip()): 
                        current_event['speakers'].append(nxt.strip()); j += 1
                    else: break
                i = j; continue
            i += 1; continue

        if not current_event['chair'] and not current_event['speakers']:
            if len(line) > 3: current_event['temp_topic'].append(line)
        i += 1

    if current_event and current_event['topic']: events.append(current_event)

    final = []
    for evt in events:
        if tr_lower(secilen_ay) in tr_lower(evt['date']):
            clean_sp = [s.replace("Vaka Sunumları:", "").strip() for s in evt['speakers'] if s.replace("Vaka Sunumları:", "").strip()]
            evt['speakers'] = clean_sp
            final.append(evt)
            
    try: final.sort(key=lambda x: int(x['date'].split()[0]))
    except: pass
    return final

def slayt_sil(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    if index < len(slides): xml_slides.remove(slides[index])

# --- WEB ARAYÜZÜ ---

# 1. Dosyaları Kontrol Et
pdf_files, sablon_path = dosyalari_bul()

if not pdf_files:
    st.error("❌ HATA: Depoda PDF dosyası bulunamadı.")
    st.stop()
if not sablon_path:
    st.error("❌ HATA: Depoda 'sablon.pptx' dosyası bulunamadı.")
    st.stop()

st.success(f"✅ Sistem Hazır: {len(pdf_files)} PDF ve Şablon algılandı.")

# 2. Ay Seçimi
secilen_ay = st.selectbox("Hangi Ayın Sunumu Hazırlanacak?", 
             ["Eylül", "Ekim", "Kasım", "Aralık", "Ocak", "Şubat", "Mart", "Nisan", "Mayıs", "Haziran"])

# 3. Çalıştırma Butonu
if st.button("🚀 Sunumu Oluştur"):
    with st.spinner('Veriler taranıyor ve slaytlar işleniyor...'):
        
        # Tüm PDF'leri Tara
        all_events = []
        for pdf in pdf_files:
            events = verileri_ayikla(pdf, secilen_ay)
            all_events.extend(events)
        
        if not all_events:
            st.warning(f"⚠️ '{secilen_ay}' ayı için programda uygun veri bulunamadı.")
        else:
            try:
                all_events.sort(key=lambda x: int(x['date'].split()[0]))
            except: pass
            
            prs = Presentation(sablon_path)
            
            # İşleme Motoru
            for event in all_events:
                num = len(event['speakers'])
                # Şablon Seçimi (1, 2, 3+ Kişi)
                t_idx = 0 if num == 1 else (1 if num == 2 else (2 if len(prs.slides)>2 else 1))
                
                source = prs.slides[t_idx]
                dest = prs.slides.add_slide(source.slide_layout)
                
                for sp in list(dest.shapes): dest.shapes._spTree.remove(sp.element)
                for shp in source.shapes:
                    newel = copy.deepcopy(shp.element)
                    dest.shapes._spTree.insert_element_before(newel, 'p:extLst')
                
                for shape in dest.shapes:
                    if not shape.has_text_frame: continue
                    for p in shape.text_frame.paragraphs:
                        txt = p.text.strip()
                        if txt == "tarih": replace_text_preserve_style(p, event['date'])
                        elif txt == "Konu adı": replace_text_preserve_style(p, event['topic'])
                        elif txt == "Başkan adı": replace_text_preserve_style(p, event['chair'])
                        
                        spk = event['speakers']
                        if t_idx == 0:
                            if (txt == "Konuşmacı adı") or ("KONUŞMACI" in shape.text_frame.text and txt == "Başkan adı"):
                                if len(spk)>0: replace_text_preserve_style(p, spk[0])
                        elif t_idx == 1:
                            if txt == "1. Konuşmacı adı" and len(spk)>0: replace_text_preserve_style(p, spk[0])
                            elif txt == "2. Konuşmacı adı" and len(spk)>1: replace_text_preserve_style(p, spk[1])
                        elif t_idx == 2:
                            if txt == "1. Konuşmacı adı" and len(spk)>0: replace_text_preserve_style(p, spk[0])
                            elif txt == "2. Konuşmacı adı" and len(spk)>1: replace_text_preserve_style(p, spk[1])
                            elif txt == "3. Konuşmacı adı" and len(spk)>2: replace_text_preserve_style(p, spk[2])
                            elif txt == "4. Konuşmacı adı" and len(spk)>3: replace_text_preserve_style(p, spk[3])

            # Şablonları Sil
            sablon_sayisi = len(prs.slides) - len(all_events)
            if len(all_events) > 0 and sablon_sayisi > 0:
                for _ in range(sablon_sayisi): slayt_sil(prs, 0)
            
            # Kaydet ve İndirme Butonu Oluştur
            binary_output = BytesIO()
            prs.save(binary_output)
            binary_output.seek(0)
            
            st.success(f"🎉 İşlem Tamam! {len(all_events)} slayt oluşturuldu.")
            
            dosya_adi = f"MSEP_{secilen_ay}_Sunumu.pptx"
            st.download_button(
                label="📥 Sunumu İndir",
                data=binary_output,
                file_name=dosya_adi,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )